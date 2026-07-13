from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches
from PIL import Image


ROOT = Path(__file__).resolve().parents[1]
ASSET_DIR = ROOT / "public" / "design-assets" / "multi-agent-worktree"
OUTPUT = ASSET_DIR / "multi-agent-worktree.pptx"

SLIDES = [
    ASSET_DIR / "slides" / "slide-worktree-overview.png",
    ASSET_DIR / "slides" / "slide-worktree-flow.png",
    ASSET_DIR / "slides" / "slide-worktree-rules.png",
]

COVERS = [
    ASSET_DIR / "covers" / "cover-worktree-main.png",
    ASSET_DIR / "covers" / "cover-worktree-agents.png",
    ASSET_DIR / "covers" / "cover-worktree-speed.png",
]


def add_contained_image(slide, image_path, x, y, width, height):
    with Image.open(image_path) as image:
        image_ratio = image.width / image.height
    box_ratio = width / height

    if image_ratio > box_ratio:
        rendered_width = width
        rendered_height = width / image_ratio
        rendered_x = x
        rendered_y = y + (height - rendered_height) / 2
    else:
        rendered_height = height
        rendered_width = height * image_ratio
        rendered_x = x + (width - rendered_width) / 2
        rendered_y = y

    slide.shapes.add_picture(
        str(image_path),
        Inches(rendered_x),
        Inches(rendered_y),
        Inches(rendered_width),
        Inches(rendered_height),
    )


def build_presentation():
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    blank = presentation.slide_layouts[6]

    for image_path in SLIDES:
        slide = presentation.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(248, 244, 232)
        add_contained_image(slide, image_path, 0, 0, 13.333, 7.5)

    cover_slide = presentation.slides.add_slide(blank)
    cover_slide.background.fill.solid()
    cover_slide.background.fill.fore_color.rgb = RGBColor(248, 244, 232)

    cover_width = 3.75
    gap = 0.45
    start_x = (13.333 - cover_width * 3 - gap * 2) / 2
    for index, image_path in enumerate(COVERS):
        add_contained_image(
            cover_slide,
            image_path,
            start_x + index * (cover_width + gap),
            0.35,
            cover_width,
            6.8,
        )

    presentation.save(OUTPUT)
    print(f"Generated {OUTPUT}")


if __name__ == "__main__":
    build_presentation()
